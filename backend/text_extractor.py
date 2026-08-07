"""
统一文字提取（通用，可被任何接口复用）。

按文件类型分发：
  .pdf            → ocr_service.process_file（自动判断文字型/图片型）
  .doc            → olefile 纯 Python 解析 OLE2 文本（零系统依赖；soffice/antiword 兜底；
                    文本过短疑扫描贴图时优先走 soffice→docx，嵌图 OCR）
  .docx           → python-docx 抽段落+表格；纯文本过短（疑扫描贴图）时 zip 内 word/media/ 嵌图 OCR
  .xlsx           → openpyxl 抽 sheet/cell 文本；纯文本过短时 xl/media/ 嵌图 OCR
  .pptx           → python-pptx 抽 slide 文本
  .png/.jpg/...   → ocr_service.extract_image_file

返回统一格式：
  {
    "text": "全文",
    "source": "pdf_text|pdf_ocr|image_ocr|docx_text",
    "page_count": int,
    "char_count": int,
  }
"""

import os
import re
import shutil
import subprocess
import asyncio
from typing import Optional

import ocr_service
import llm_service
import event_service


# 大文件 early-exit 阈值:扫描版 PDF 总页数 > 此值才启动 LLM 初判 + 采样
# 小文件全文 OCR 也就几十秒,LLM 初判反而不划算
OCR_EARLY_EXIT_THRESHOLD = int(os.getenv("OCR_EARLY_EXIT_THRESHOLD", "10"))


_DOCX_EXT = ".docx"
_XLS_EXT = ".xls"
_XLSX_EXT = ".xlsx"
_PPTX_EXT = ".pptx"
_PDF_EXT = ".pdf"
_GIF_EXT = ".gif"
_IMAGE_EXTS = {".png", ".jpg", ".jpeg", ".bmp", ".tiff", ".webp"}

# Office 嵌图 OCR(扫描件贴进 Word/Excel 的场景):纯文本短于阈值才触发,
# 正常文档里的 logo/公章图 OCR 只产噪音且占全局 OCR 锁,故不无条件跑
OFFICE_IMG_OCR_TEXT_THRESHOLD = 80
OFFICE_IMG_OCR_MAX_IMAGES = 10    # 单文档最多 OCR 的嵌图数
OFFICE_IMG_OCR_MIN_SIDE = 150     # 图片短边小于此像素视为图标/装饰图,跳过
_ZIP_MEDIA_RASTER_EXTS = {".png", ".jpg", ".jpeg", ".bmp", ".tiff", ".webp"}


# NUL/C0 控制字符清洗表:pdfplumber/python-docx/OCR 偶发混入 0x00 等控制符,
# PostgreSQL text 列不接受 NUL,下游 LLM 也无需这些字符。在抽取出口统一剔除
# (保留 \t \n \r),作为 CRUD 层清洗之外的源头防御。
_CTRL_DELETE = {c: None for c in range(0x20) if c not in (0x09, 0x0A, 0x0D)}
_CTRL_DELETE[0x7F] = None  # DEL


def _sanitize_text(text: Optional[str]) -> Optional[str]:
    """去除 NUL/C0 控制字符(保留 \t\n\r)和 DEL。None/空串原样返回。"""
    if not text:
        return text
    return text.translate(_CTRL_DELETE)


def _sanitize_extracted(result: dict) -> dict:
    """在抽取出口清洗 text 字段并同步 char_count,确保下游拿到的是干净文本。"""
    if result and result.get("text"):
        cleaned = _sanitize_text(result["text"])
        result["text"] = cleaned
        if "char_count" in result:
            result["char_count"] = len(cleaned) if cleaned else 0
    return result


_ANTIWORD_CACHE = {"path": None, "checked": False}
_SOFFICE_CACHE = {"path": None, "checked": False}

# soffice 常见安装路径(与 template_service 对齐)
_SOFFICE_CANDIDATE_PATHS = [
    r"C:\Program Files\LibreOffice\program\soffice.exe",
    r"C:\Program Files (x86)\LibreOffice\program\soffice.exe",
    "/usr/bin/soffice",
    "/usr/bin/libreoffice",
    "/Applications/LibreOffice.app/Contents/MacOS/soffice",
]


def _find_antiword() -> Optional[str]:
    """查找 antiword 可执行文件路径(懒缓存)。"""
    if not _ANTIWORD_CACHE["checked"]:
        _ANTIWORD_CACHE["path"] = shutil.which("antiword")
        _ANTIWORD_CACHE["checked"] = True
    return _ANTIWORD_CACHE["path"]


def _find_soffice() -> Optional[str]:
    """查找 LibreOffice soffice 可执行文件路径(懒缓存)。"""
    if _SOFFICE_CACHE["checked"]:
        return _SOFFICE_CACHE["path"]
    for name in ("soffice", "soffice.exe", "libreoffice", "libreoffice.exe"):
        found = shutil.which(name)
        if found:
            _SOFFICE_CACHE["path"] = found
            break
    if not _SOFFICE_CACHE["path"]:
        for cand in _SOFFICE_CANDIDATE_PATHS:
            if os.path.exists(cand):
                _SOFFICE_CACHE["path"] = cand
                break
    _SOFFICE_CACHE["checked"] = True
    return _SOFFICE_CACHE["path"]


def _extract_doc_via_soffice(soffice: str, file_path: str) -> dict:
    """用 LibreOffice 把旧版 .doc 转成 .docx 到临时目录,再复用 _extract_docx 抽文本。

    转 docx(而非 txt)能保留表格结构,和 .docx 抽取逻辑一致。
    """
    import tempfile
    src = os.path.abspath(file_path)
    with tempfile.TemporaryDirectory(prefix="doc2docx_") as tmpdir:
        cmd = [
            soffice, "--headless",
            "--convert-to", "docx:MS Word 2007 XML",
            "--outdir", tmpdir, src,
        ]
        try:
            subprocess.run(cmd, timeout=90, check=True, capture_output=True)
        except subprocess.TimeoutExpired:
            raise ValueError(".doc 转换超时(文件可能过大或损坏)")
        except subprocess.CalledProcessError as e:
            err = (e.stderr or b"").decode("utf-8", errors="replace")[:200]
            raise ValueError(f".doc 转换失败(可能已损坏或非标准 .doc 格式): {err}")
        out_docx = os.path.join(
            tmpdir, os.path.splitext(os.path.basename(src))[0] + ".docx"
        )
        if not os.path.exists(out_docx):
            raise ValueError(".doc 转换未生成 docx(LibreOffice 可能无法识别该文件)")
        result = _extract_docx(out_docx)
    # 保留嵌图 OCR 标记:docx_text→doc_text,docx_text+img_ocr→doc_text+img_ocr,docx_img_ocr→doc_img_ocr
    result["source"] = (result.get("source") or "docx_text").replace("docx", "doc", 1)
    return result


def office_to_pdf(src_path: str, dst_dir: str) -> str:
    """soffice 把 Office 文档(doc/docx/xls/xlsx/ppt/pptx)转 PDF 到 dst_dir。

    返回生成的 PDF 绝对路径(文件名=源 basename+.pdf)。
    soffice 缺失/转换失败抛 RuntimeError。同步阻塞,调用方需 asyncio.to_thread。
    """
    soffice = _find_soffice()
    if not soffice:
        raise RuntimeError("服务器未安装 LibreOffice,无法预览 Office 原件")
    cmd = [
        soffice, "--headless",
        "--convert-to", "pdf",
        "--outdir", os.path.abspath(dst_dir),
        os.path.abspath(src_path),
    ]
    try:
        subprocess.run(cmd, timeout=90, check=True, capture_output=True)
    except subprocess.TimeoutExpired:
        raise RuntimeError("Office 转 PDF 超时(文件可能过大或损坏)")
    except subprocess.CalledProcessError as e:
        err = (e.stderr or b"").decode("utf-8", errors="replace")[:200]
        raise RuntimeError(f"Office 转 PDF 失败: {err}")
    out_pdf = os.path.join(
        os.path.abspath(dst_dir),
        os.path.splitext(os.path.basename(src_path))[0] + ".pdf")
    if not os.path.exists(out_pdf):
        raise RuntimeError("Office 转 PDF 未生成输出文件(LibreOffice 可能无法识别该文件)")
    return out_pdf


def _extract_doc_via_olefile(file_path: str) -> dict:
    """纯 Python 解析旧版 .doc(OLE2 复合文档)文本,零系统依赖。

    按 [MS-DOC] 规范:读 WordDocument 流的 FIB 定位 piece table(在 Table 流里),
    逐片按 fCompressed 标志位解码(0=UTF-16LE, 1=CP1252 8-bit)。
    这是 antiword/textract 用的标准算法,比"扫可读字符"干净得多。
    """
    import struct
    import olefile

    if not olefile.isOleFile(file_path):
        raise ValueError(".doc 不是有效的 OLE2 文件(可能已损坏或实为其他格式)")

    ole = olefile.OleFileIO(file_path)
    try:
        if not ole.exists("WordDocument"):
            raise ValueError(".doc 缺少 WordDocument 流(非标准 Word 文档)")
        with ole.openstream("WordDocument") as s:
            doc = s.read()

        # FIB: fComplex/fWhichTblStm 在 base 的 flags(offset 10, 2字节 little-endian)
        # bit 9 (0x0200) = fWhichTblStm: 1→用 1Table, 0→用 0Table
        flags = struct.unpack_from("<H", doc, 0x000A)[0]
        table_name = "1Table" if (flags & 0x0200) else "0Table"

        # fcClx / lcbClx 在 FibRgFcLcb97:fcMin=..., fcClx 位于 offset 0x01A2, lcbClx 0x01A6
        fc_clx = struct.unpack_from("<I", doc, 0x01A2)[0]
        lcb_clx = struct.unpack_from("<I", doc, 0x01A6)[0]
        # ccpText: 正文字符数,FIB offset 0x004C
        ccp_text = struct.unpack_from("<I", doc, 0x004C)[0]

        text = ""
        if ole.exists(table_name) and lcb_clx > 0:
            with ole.openstream(table_name) as ts:
                table = ts.read()
            clx = table[fc_clx:fc_clx + lcb_clx]
            text = _parse_doc_piece_table(doc, clx, ccp_text)

        # 兜底:piece table 解析不出内容时,直接从 fcMin 起按 CP1252 粗提
        if not text.strip():
            text = _doc_fallback_scan(doc, ccp_text)
    finally:
        ole.close()

    text = text.strip()
    if not text:
        raise ValueError(".doc 解析后无文字(可能是扫描件转存的图片型 doc,或加密文档)")
    return {
        "text": text,
        "source": "doc_text",
        "page_count": 1,
        "char_count": len(text),
    }


def _parse_doc_piece_table(doc: bytes, clx: bytes, ccp_text: int) -> str:
    """从 Clx(复杂格式区)解析 piece table,拼出正文文本。"""
    import struct

    # Clx = 可选的 RgPrc(0x01 开头) + Pcdt(0x02 开头)
    i = 0
    n = len(clx)
    while i < n and clx[i] == 0x01:
        # Prc: 0x01 + cbGrpprl(2字节) + grpprl
        if i + 3 > n:
            break
        cb = struct.unpack_from("<H", clx, i + 1)[0]
        i += 3 + cb
    if i >= n or clx[i] != 0x02:
        return ""
    # Pcdt: 0x02 + lcb(4字节) + PlcPcd
    lcb = struct.unpack_from("<I", clx, i + 1)[0]
    plc = clx[i + 5: i + 5 + lcb]

    # PlcPcd = (n+1) 个 CP(4字节) + n 个 Pcd(8字节);由此反推 n
    # lcb = 4*(n+1) + 8*n = 12n + 4  → n = (lcb-4)/12
    if lcb < 4:
        return ""
    n_pieces = (lcb - 4) // 12
    cps = [struct.unpack_from("<I", plc, k * 4)[0] for k in range(n_pieces + 1)]
    pcd_base = (n_pieces + 1) * 4

    parts = []
    for k in range(n_pieces):
        pcd = plc[pcd_base + k * 8: pcd_base + k * 8 + 8]
        fc_field = struct.unpack_from("<I", pcd, 2)[0]
        # fc 的 bit30 = fCompressed:1→CP1252 8-bit(fc>>1 为字节偏移),0→UTF-16LE
        f_compressed = (fc_field & 0x40000000) != 0
        fc = fc_field & 0x3FFFFFFF
        cp_start, cp_end = cps[k], cps[k + 1]
        char_cnt = cp_end - cp_start
        if char_cnt <= 0:
            continue
        if f_compressed:
            off = fc // 2
            raw = doc[off: off + char_cnt]
            parts.append(raw.decode("cp1252", errors="replace"))
        else:
            off = fc
            raw = doc[off: off + char_cnt * 2]
            parts.append(raw.decode("utf-16-le", errors="replace"))

    text = "".join(parts)
    return _clean_doc_text(text)


def _doc_fallback_scan(doc: bytes, ccp_text: int) -> str:
    """piece table 不可用时的兜底:从 fcMin 起按 UTF-16/CP1252 粗解。"""
    import struct
    fc_min = struct.unpack_from("<I", doc, 0x0018)[0]
    chunk = doc[fc_min: fc_min + max(ccp_text, 0) * 2] if ccp_text else doc[fc_min:]
    # 先试 UTF-16LE,失败再 CP1252
    try:
        text = chunk.decode("utf-16-le", errors="ignore")
    except Exception:
        text = chunk.decode("cp1252", errors="ignore")
    return _clean_doc_text(text)


def _clean_doc_text(text: str) -> str:
    """清洗 .doc 提取出的控制字符:Word 特殊码 → 换行/制表/空;删其余控制符。"""
    if not text:
        return ""
    # Word 特殊字符映射
    trans = {
        0x0007: "\n",   # 单元格/行结束
        0x000D: "\n",   # 段落结束
        0x000B: "\n",   # 手动换行
        0x000C: "\n",   # 分页
        0x001E: "-",    # 不间断连字符
        0x001F: "",     # 可选连字符
        0x00A0: " ",    # 不间断空格
        0x2002: " ", 0x2003: " ", 0x2009: " ",
        0xFC01: "", 0xF020: "",
    }
    out = []
    for ch in text:
        o = ord(ch)
        if o in trans:
            out.append(trans[o])
        elif o == 0x09 or o == 0x0A:
            out.append(ch)
        elif o < 0x20:
            continue          # 其余控制字符丢弃
        elif o == 0xFFFF or o == 0xFFFE:
            continue
        else:
            out.append(ch)
    # 规整多余空行
    cleaned = "\n".join(line.rstrip() for line in "".join(out).split("\n"))
    while "\n\n\n" in cleaned:
        cleaned = cleaned.replace("\n\n\n", "\n\n")
    return cleaned.strip()


def _extract_doc(file_path: str) -> dict:
    """旧版 .doc 二进制文本抽取。

    优先用纯 Python 的 olefile 解析(零系统依赖);olefile 成功但文本极短(疑扫描贴图)
    时优先改走 soffice→docx 路线以触发嵌图 OCR;olefile 失败再尝试 soffice / antiword;
    都不行才抛 ValueError。
    """
    # 1) 纯 Python olefile(首选,服务器无需装任何系统包)
    last_err = None
    ole_result = None
    try:
        ole_result = _extract_doc_via_olefile(file_path)
    except Exception as e:
        last_err = e

    soffice = _find_soffice()
    if ole_result is not None:
        # olefile 成功但文本极短:可能是扫描件贴图型 .doc(OLE 嵌图纯 Python 拆不动),
        # 有 soffice 时改走 soffice→docx 路线(_extract_docx 会对嵌图 OCR);失败保留 olefile 结果
        if len(ole_result.get("text") or "") < OFFICE_IMG_OCR_TEXT_THRESHOLD and soffice:
            try:
                return _extract_doc_via_soffice(soffice, file_path)
            except Exception:
                pass
        return ole_result

    # 2) soffice 兜底(装了 LibreOffice 的环境)
    if soffice:
        try:
            return _extract_doc_via_soffice(soffice, file_path)
        except Exception as e:
            last_err = e

    # 3) antiword 兜底(装了 antiword 的环境,如本地 Windows)
    antiword = _find_antiword()
    if antiword:
        try:
            result = subprocess.run(
                [antiword, "-m", "UTF-8.txt", os.path.abspath(file_path)],
                timeout=60, check=True, capture_output=True,
            )
            text = result.stdout.decode("utf-8", errors="replace").strip()
            if text.strip():
                return {"text": text, "source": "doc_text",
                        "page_count": 1, "char_count": len(text)}
        except Exception as e:
            last_err = e

    raise ValueError(f".doc 解析失败(文件可能已损坏、加密或为图片型): {last_err}")


def _ocr_zip_media(file_path: str, media_prefix: str) -> tuple:
    """OOXML(zip)文档 media_prefix(word/media/ 或 xl/media/)下嵌图 OCR。

    返回 (ocr_text, ocr_image_count)。只处理光栅图(跳过 emf/wmf/gif 等);
    按条目名排序保证确定性;OCR 满 OFFICE_IMG_OCR_MAX_IMAGES 张即停;
    短边 < OFFICE_IMG_OCR_MIN_SIDE 的小图(图标/logo)跳过;单图异常跳过不杀整体。
    """
    import io
    import tempfile
    import zipfile
    from PIL import Image

    try:
        zf = zipfile.ZipFile(file_path)
    except Exception:
        return "", 0
    texts: list[str] = []
    ocr_count = 0
    try:
        names = sorted(
            n for n in zf.namelist()
            if n.startswith(media_prefix)
            and os.path.splitext(n)[1].lower() in _ZIP_MEDIA_RASTER_EXTS
        )
        for name in names:
            if ocr_count >= OFFICE_IMG_OCR_MAX_IMAGES:
                break
            tmp_path = None
            try:
                data = zf.read(name)
                with Image.open(io.BytesIO(data)) as im:
                    if min(im.size) < OFFICE_IMG_OCR_MIN_SIDE:
                        continue
                with tempfile.NamedTemporaryFile(
                        delete=False, suffix=os.path.splitext(name)[1].lower()) as tmp:
                    tmp.write(data)
                    tmp_path = tmp.name
                ocr_result = ocr_service.run_ocr(tmp_path, cls=True)
                ocr_count += 1
                if ocr_result and ocr_result[0]:
                    for line in ocr_result[0]:
                        if float(line[1][1]) > 0.3:
                            texts.append(line[1][0])
            except Exception as e:
                print(f"[text_extractor] 嵌图 OCR 失败(跳过): {name}: {e}")
            finally:
                if tmp_path:
                    try:
                        os.remove(tmp_path)
                    except OSError:
                        pass
    finally:
        zf.close()
    return "\n".join(texts).strip(), ocr_count


def _log_office_img_ocr(file_path: str, media_prefix: str, ocr_images: int, ocr_chars: int) -> None:
    """嵌图 OCR 触发事件(事件流可观测命中率);失败静默,不影响抽取。"""
    try:
        event_service.log_event(
            event_service.INFO,
            event_service.CATEGORY_FILE_OCR_SAMPLED,
            f"Office 文档嵌入图片 OCR:{os.path.basename(file_path)}",
            context={
                "filename": os.path.basename(file_path),
                "media_prefix": media_prefix,
                "ocr_images": ocr_images,
                "ocr_chars": ocr_chars,
            },
        )
    except Exception:
        pass


def _merge_office_img_ocr(file_path: str, text: str, source: str, media_prefix: str) -> tuple:
    """纯文本过短时对 OOXML 嵌图 OCR 并合并文本。返回 (text, source)。"""
    if len(text) >= OFFICE_IMG_OCR_TEXT_THRESHOLD:
        return text, source
    ocr_text, ocr_images = _ocr_zip_media(file_path, media_prefix)
    if not ocr_text:
        return text, source
    if text:
        text = f"{text}\n\n--- 嵌入图片 OCR ---\n{ocr_text}"
        source = f"{source}+img_ocr"
    else:
        text = ocr_text
        source = source.replace("_text", "_img_ocr", 1)
    _log_office_img_ocr(file_path, media_prefix, ocr_images, len(ocr_text))
    return text, source


def _extract_docx(file_path: str) -> dict:
    """python-docx 抽 docx 全文：段落 + 表格 cell 文字。"""
    from docx import Document
    doc = Document(file_path)

    parts: list[str] = []

    # 段落
    for para in doc.paragraphs:
        t = para.text.strip()
        if t:
            parts.append(t)

    # 表格
    for tbl_idx, table in enumerate(doc.tables, 1):
        parts.append(f"--- 表 {tbl_idx} ---")
        for row in table.rows:
            cells = [cell.text.strip() for cell in row.cells]
            cells = [c for c in cells if c]
            if cells:
                parts.append(" | ".join(cells))

    text = "\n".join(parts).strip()
    # 纯文本过短:疑扫描件贴图文档(证件扫描贴进 Word),对 word/media/ 嵌图 OCR
    text, source = _merge_office_img_ocr(file_path, text, "docx_text", "word/media/")
    return {
        "text": text,
        "source": source,
        "page_count": 1,                 # docx 没有"页"概念，记 1
        "char_count": len(text),
    }


def _extract_xlsx(file_path: str) -> dict:
    """openpyxl 抽 xlsx 文本：逐 sheet 读取非空 cell。"""
    from openpyxl import load_workbook

    wb = load_workbook(file_path, read_only=True, data_only=True)
    parts: list[str] = []
    max_lines = 5000  # 防止超大表格抽取过多无效文本；LLM 层还会二次截断

    sheet_count = len(wb.sheetnames)
    try:
        for ws in wb.worksheets:
            if len(parts) >= max_lines:
                break
            parts.append(f"--- Sheet: {ws.title} ---")
            for row in ws.iter_rows():
                cells = []
                for cell in row:
                    if cell.value is None:
                        continue
                    value = str(cell.value).strip()
                    if value:
                        cells.append(f"{cell.coordinate}: {value}")
                if cells:
                    parts.append(" | ".join(cells))
                if len(parts) >= max_lines:
                    parts.append("...[表格内容过长，已截断]...")
                    break
    finally:
        wb.close()

    text = "\n".join(parts).strip()
    # 纯文本过短:疑扫描件贴图表格,对 xl/media/ 嵌图 OCR
    text, source = _merge_office_img_ocr(file_path, text, "xlsx_text", "xl/media/")
    return {
        "text": text,
        "source": source,
        "page_count": sheet_count,
        "char_count": len(text),
    }


def _extract_xls(file_path: str) -> dict:
    """xlrd 抽旧版 xls 文本：逐 sheet 读取非空 cell。"""
    import xlrd

    book = xlrd.open_workbook(file_path)
    parts: list[str] = []
    max_lines = 5000

    for sheet in book.sheets():
        if len(parts) >= max_lines:
            break
        parts.append(f"--- Sheet: {sheet.name} ---")
        for r in range(sheet.nrows):
            cells = []
            for c in range(sheet.ncols):
                value = sheet.cell_value(r, c)
                if value is None:
                    continue
                value = str(value).strip()
                if value:
                    # xlrd 用 0-based,展示成 R/C 避免复杂列号转换
                    cells.append(f"R{r + 1}C{c + 1}: {value}")
            if cells:
                parts.append(" | ".join(cells))
            if len(parts) >= max_lines:
                parts.append("...[表格内容过长，已截断]...")
                break

    text = "\n".join(parts).strip()
    return {
        "text": text,
        "source": "xls_text",
        "page_count": book.nsheets,
        "char_count": len(text),
    }
def _extract_pptx(file_path: str) -> dict:
    """python-pptx 抽 pptx 文本：逐 slide 抽 shape.text。"""
    from pptx import Presentation

    prs = Presentation(file_path)
    parts: list[str] = []
    for idx, slide in enumerate(prs.slides, 1):
        slide_lines = []
        for shape in slide.shapes:
            text = getattr(shape, "text", "")
            if text and text.strip():
                slide_lines.append(text.strip())
        if slide_lines:
            parts.append(f"--- Slide {idx} ---")
            parts.extend(slide_lines)

    text = "\n".join(parts).strip()
    # 纯文本过短:疑扫描件贴图 PPT,对 ppt/media/ 嵌图 OCR(与 docx/xlsx 同口径)
    text, source = _merge_office_img_ocr(file_path, text, "pptx_text", "ppt/media/")
    return {
        "text": text,
        "source": source,
        "page_count": len(prs.slides),
        "char_count": len(text),
    }


def _cleanup_ocr_dir(task_id: str) -> None:
    """删除 OCR 渲染中间产物 output/{task_id}/。
    业务审核/摘要场景下,PNG 只是 OCR 中间产物,文字抽完即可丢弃。
    失败不抛(磁盘清理是 best-effort)。
    """
    try:
        import shutil
        d = os.path.join(ocr_service.OUTPUT_DIR, task_id)
        if os.path.isdir(d):
            shutil.rmtree(d, ignore_errors=True)
    except Exception as e:
        print(f"[text_extractor] 清理 {task_id} 失败(忽略): {e}")


def _ocr_single_page(file_path: str, task_id: str, page_index_0based: int,
                     render_dpi: int = 0) -> str:
    """只取指定页(0-based 索引)的文本。用于大文件 early-exit 时抓末页盖章/合计。

    逐页混合:该页含大图 → 渲染+OCR;无大图 → 直接读文本层(数字页不浪费 OCR)。
    内联实现,不污染 ocr_service 接口。复用 ocr_service.run_ocr 引擎锁,跨线程池安全。
    """
    import pdfplumber
    import pypdfium2

    # 数字页:无大图 → 文本层直接返回;文本层为空(文字转曲/伪数字页)继续走 OCR
    with pdfplumber.open(file_path) as plumber:
        ppage = plumber.pages[page_index_0based]
        if not ocr_service._page_has_big_image(ppage):
            text = (ppage.extract_text() or "").strip()
            if text:
                return text

    # 扫描页:渲染 + OCR
    img_dir = os.path.join(ocr_service.OUTPUT_DIR, task_id, "images")
    os.makedirs(img_dir, exist_ok=True)

    pdf = pypdfium2.PdfDocument(file_path)
    try:
        page = pdf[page_index_0based]
        bitmap = page.render(scale=ocr_service._render_scale(render_dpi))
        pil_image = bitmap.to_pil()
        pil_image, _ = ocr_service._downscale_if_too_large(pil_image)
        img_filename = f"page_{page_index_0based + 1}.png"
        img_path = os.path.join(img_dir, img_filename)
        pil_image.save(img_path, "PNG")
    finally:
        pdf.close()

    # OCR 这一张图
    ocr_result = ocr_service.run_ocr(img_path, cls=True)
    lines = []
    if ocr_result and ocr_result[0]:
        for line in ocr_result[0]:
            text = line[1][0]
            conf = float(line[1][1])
            if conf > 0.3:
                lines.append(text)
    return "\n".join(lines)


def _extract_pdf(file_path: str, render_dpi: int = 0) -> dict:
    """PDF 文字提取,带大文件 early-exit 优化。

    render_dpi>0 时按该 DPI 渲染扫描页(默认 200;画像管线传 300 提升小字识别)。

    流程:
    1. 文字型 PDF(有文字层) → pdfplumber 全文(秒级)
    2. 扫描版 ≤ OCR_EARLY_EXIT_THRESHOLD 页 → 走原 OCR 全文
    3. 扫描版 > 阈值页 → 先 OCR 前 2 页 → LLM 初判
        - 是大表类(流水/社保/证券): 再 OCR 末页 = 3 页采样返回
        - 否则: OCR 第 3 页起的剩余页 = 全文返回
        - LLM 抽风: 也走采样(激进保速度)
    """
    import uuid
    base = os.path.splitext(os.path.basename(file_path))[0]
    task_id = f"fetched_{base}_{uuid.uuid4().hex[:8]}"

    try:
        # === 1. 文字型 PDF ===
        pdf_type = ocr_service.detect_pdf_type(file_path)
        if pdf_type == "text":
            try:
                pages = ocr_service.extract_text_pdf(file_path)
                # 文字型判定是文档级的,逐页补漏:文本层为空的页(文字转曲伪数字页/
                # 追加的无文本层扫描页)单页 OCR 兜底,不静默吞内容
                ocr_filled = False
                for i, p in enumerate(pages):
                    if (p.get("text") or "").strip():
                        continue
                    try:
                        ocr_text = _ocr_single_page(file_path, task_id, i,
                                                    render_dpi=render_dpi)
                    except Exception as e:
                        print(f"[text_extractor] 第 {i + 1} 页兜底 OCR 失败(留空): {e}")
                        continue
                    if ocr_text.strip():
                        p["text"] = ocr_text
                        ocr_filled = True
                text = "\n\n".join(p.get("text", "") for p in pages).strip()
                return {
                    "text": text,
                    "source": "pdf_text+ocr" if ocr_filled else "pdf_text",
                    "page_count": len(pages),
                    "char_count": len(text),
                }
            except Exception as e:
                # 文本层损坏(畸形 PDF 触发 pdfminer 内部 IndexError 等) → 按扫描件走 OCR 兜底
                print(f"[text_extractor] PDF 文本层解析失败,降级按扫描件 OCR: {type(e).__name__}: {e}")

        # === 2. 拿总页数 ===
        import pypdfium2
        pdf = pypdfium2.PdfDocument(file_path)
        try:
            total_pages = len(pdf)
        finally:
            pdf.close()

        # === 3. 小文件直接全 OCR(逐页混合:扫描页 OCR,数字页读文本层) ===
        if total_pages <= OCR_EARLY_EXIT_THRESHOLD:
            pages = ocr_service.extract_mixed_pdf(file_path, task_id, max_ocr_pages=0,
                                                  render_dpi=render_dpi)
            text = "\n\n".join(p.get("text", "") for p in pages).strip()
            return {
                "text": text,
                "source": "pdf_ocr",
                "page_count": len(pages),
                "char_count": len(text),
            }

        # === 4. 大文件 early-exit:先取前 2 页(逐页混合) ===
        head_pages = ocr_service.extract_mixed_pdf(file_path, task_id, max_ocr_pages=2,
                                                   render_dpi=render_dpi)
        head_text = "\n\n".join(p.get("text", "") for p in head_pages[:2]).strip()

        # === 5. LLM 初判 ===
        # task_id 传短 id(上面构造的 fetched_*),不是文件全路径——路径会撑爆
        # ai_api_calls.task_id varchar(64) 且日志里没有检索价值
        verdict = llm_service.detect_large_table_doc(head_text, task_id=task_id)
        is_large_table = bool(verdict.get("is_large_table"))
        doc_type = verdict.get("doc_type") or "unknown"
        is_fallback = bool(verdict.get("_fallback"))

        # === 6a. 是大表类(或 LLM 抽风) → 加一页末页就返回 ===
        if is_large_table:
            try:
                tail_text = _ocr_single_page(file_path, task_id, total_pages - 1,
                                             render_dpi=render_dpi)
            except Exception as e:
                print(f"[text_extractor] 末页 OCR 失败,只用前 2 页: {e}")
                tail_text = ""

            sampled_text = head_text
            if tail_text:
                sampled_text += "\n\n--- 中间页未识别 ---\n\n" + tail_text
            sampled_text += (
                f"\n\n[已采样 OCR: 共 {total_pages} 页, 实际识别第 1,2,{total_pages} 页. "
                f"判定为 {doc_type}{'(LLM 降级)' if is_fallback else ''}]"
            )

            # 记一条事件,事件流可观测采样命中率
            try:
                event_service.log_event(
                    event_service.INFO,
                    event_service.CATEGORY_FILE_OCR_SAMPLED,
                    f"文件采样 OCR:共 {total_pages} 页,识别 3 页({doc_type})",
                    context={
                        "filename": os.path.basename(file_path),
                        "total_pages": total_pages,
                        "sampled_pages": [1, 2, total_pages],
                        "doc_type": doc_type,
                        "llm_fallback": is_fallback,
                        "confidence": verdict.get("confidence", 0),
                    },
                )
            except Exception:
                pass

            return {
                "text": sampled_text,
                "source": "pdf_ocr_sampled",
                "page_count": total_pages,
                "char_count": len(sampled_text),
            }

        # === 6b. 不是大表类 → 取全文(逐页混合;前 2 页缓存重复跑,代价小,逻辑简单) ===
        pages = ocr_service.extract_mixed_pdf(file_path, task_id, max_ocr_pages=0,
                                              render_dpi=render_dpi)
        text = "\n\n".join(p.get("text", "") for p in pages).strip()
        return {
            "text": text,
            "source": "pdf_ocr",
            "page_count": len(pages),
            "char_count": len(text),
        }

    finally:
        _cleanup_ocr_dir(task_id)


def _extract_image(file_path: str) -> dict:
    """图片走 PaddleOCR。"""
    import uuid
    base = os.path.splitext(os.path.basename(file_path))[0]
    task_id = f"fetched_{base}_{uuid.uuid4().hex[:8]}"
    try:
        pages = ocr_service.extract_image_file(file_path, task_id)
        text = "\n".join(p.get("text", "") for p in pages).strip()
        return {
            "text": text,
            "source": "image_ocr",
            "page_count": len(pages),
            "char_count": len(text),
        }
    finally:
        _cleanup_ocr_dir(task_id)


def _extract_gif(file_path: str) -> dict:
    """GIF 只取第一帧转 PNG 后走 OCR。"""
    from PIL import Image
    import tempfile

    base = os.path.splitext(os.path.basename(file_path))[0]
    tmp_path = os.path.join(tempfile.gettempdir(), f"{base}_gif_first_frame.png")
    try:
        with Image.open(file_path) as im:
            im.seek(0)
            im.convert("RGB").save(tmp_path, "PNG")
        result = _extract_image(tmp_path)
        result["source"] = "gif_first_frame_ocr"
        return result
    finally:
        try:
            if os.path.exists(tmp_path):
                os.remove(tmp_path)
        except OSError:
            pass


async def extract_text(file_path: str, mime_type: Optional[str] = None,
                       render_dpi: int = 0) -> dict:
    """统一文字提取入口（异步，把同步阻塞代码扔到线程池）。

    render_dpi>0 时 PDF 扫描页按该 DPI 渲染(默认 200;画像管线传 300 提升小字识别)。

    抛出：
      ValueError - 不支持的扩展名
      FileNotFoundError
      其他异常 - 透传给调用方
    """
    if not os.path.exists(file_path):
        raise FileNotFoundError(f"文件不存在: {file_path}")

    ext = os.path.splitext(file_path)[1].lower()

    if ext == _DOCX_EXT:
        return _sanitize_extracted(await asyncio.to_thread(_extract_docx, file_path))

    if ext == _XLSX_EXT:
        return _sanitize_extracted(await asyncio.to_thread(_extract_xlsx, file_path))

    if ext == _PPTX_EXT:
        return _sanitize_extracted(await asyncio.to_thread(_extract_pptx, file_path))

    if ext == _XLS_EXT:
        return _sanitize_extracted(await asyncio.to_thread(_extract_xls, file_path))

    if ext == _GIF_EXT:
        return _sanitize_extracted(await asyncio.to_thread(_extract_gif, file_path))

    if ext == ".doc":
        return _sanitize_extracted(await asyncio.to_thread(_extract_doc, file_path))

    if ext == _PDF_EXT:
        return _sanitize_extracted(await asyncio.to_thread(_extract_pdf, file_path, render_dpi))

    if ext in _IMAGE_EXTS:
        return _sanitize_extracted(await asyncio.to_thread(_extract_image, file_path))

    raise ValueError(f"不支持的文件类型: {ext}（支持 .pdf/.doc/.docx/.xls/.xlsx/.pptx/.gif/{'/'.join(sorted(_IMAGE_EXTS))}）")


def normalize_text(text: str, max_chars: Optional[int] = None) -> str:
    """规整文本：合并多余空白、可选截断。
    给 LLM 用之前调一下，避免无意义 token 浪费。
    """
    if not text:
        return ""
    # 合并多余空白行
    text = re.sub(r"\n{3,}", "\n\n", text)
    text = re.sub(r"[ \t]+", " ", text)
    text = text.strip()

    if max_chars and len(text) > max_chars:
        head = max_chars // 2
        tail = max_chars - head
        text = text[:head] + f"\n\n...[省略 {len(text) - max_chars} 字]...\n\n" + text[-tail:]
    return text
